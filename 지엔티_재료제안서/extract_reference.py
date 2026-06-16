"""
참조 PPT(재료 정리.pptx)에서 슬라이드별 이미지·텍스트 추출.
- 이미지: _extracted/sNN-imgM.png 형태로 저장
- 텍스트: _extracted/data.json
"""

import os
import json
from pptx import Presentation
from pptx.util import Emu

SRC = r"C:\Users\Administrator\Desktop\JM_클로드\지엔티_재료제안서\재료 정리.pptx"
OUT = r"C:\Users\Administrator\Desktop\JM_클로드\지엔티_재료제안서\_extracted"
os.makedirs(OUT, exist_ok=True)

prs = Presentation(SRC)
sw, sh = prs.slide_width, prs.slide_height
print(f"Slide size: {Emu(sw).inches:.2f} x {Emu(sh).inches:.2f} inch  ({prs.slide_width}, {prs.slide_height} EMU)")

data = []

for s_idx, slide in enumerate(prs.slides, start=1):
    s_data = {"slide": s_idx, "texts": [], "images": []}
    img_count = 0

    for shape in slide.shapes:
        # 텍스트
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = "".join(run.text for run in p.runs)
                if t.strip():
                    s_data["texts"].append({
                        "text": t.strip(),
                        "x_in": round(Emu(shape.left).inches, 2) if shape.left else 0,
                        "y_in": round(Emu(shape.top).inches, 2) if shape.top else 0,
                        "w_in": round(Emu(shape.width).inches, 2) if shape.width else 0,
                    })

        # 이미지 (Picture shape)
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            img_count += 1
            img = shape.image
            ext = img.ext  # 'png', 'jpeg' etc
            fname = f"s{s_idx:02d}-img{img_count}.{ext}"
            with open(os.path.join(OUT, fname), "wb") as f:
                f.write(img.blob)
            s_data["images"].append({
                "file": fname,
                "x_in": round(Emu(shape.left).inches, 2),
                "y_in": round(Emu(shape.top).inches, 2),
                "w_in": round(Emu(shape.width).inches, 2),
                "h_in": round(Emu(shape.height).inches, 2),
            })

    data.append(s_data)
    print(f"Slide {s_idx}: {len(s_data['texts'])} text(s), {len(s_data['images'])} image(s)")

with open(os.path.join(OUT, "data.json"), "w", encoding="utf-8") as f:
    json.dump(data, f, ensure_ascii=False, indent=2)

print(f"\nDone → {OUT}\\data.json")
