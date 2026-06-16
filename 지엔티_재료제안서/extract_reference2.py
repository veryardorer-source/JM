"""갱신된 재료 정리.pptx에서 모든 슬라이드의 이미지·텍스트 재추출."""

import os
import json
from pptx import Presentation
from pptx.util import Emu

SRC = r"C:\Users\Administrator\Desktop\JM_클로드\지엔티_재료제안서\재료 정리.pptx"
OUT = r"C:\Users\Administrator\Desktop\JM_클로드\지엔티_재료제안서\_extracted2"
os.makedirs(OUT, exist_ok=True)

prs = Presentation(SRC)
print(f"Total slides: {len(prs.slides)}")

data = []
for s_idx, slide in enumerate(prs.slides, start=1):
    s_data = {"slide": s_idx, "texts": [], "images": []}
    img_n = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs).strip()
                if t:
                    s_data["texts"].append({
                        "text": t,
                        "y_in": round(Emu(shape.top).inches, 2) if shape.top else 0,
                        "x_in": round(Emu(shape.left).inches, 2) if shape.left else 0,
                    })
        if shape.shape_type == 13:  # PICTURE
            img_n += 1
            ext = shape.image.ext
            fname = f"s{s_idx:02d}-img{img_n}.{ext}"
            with open(os.path.join(OUT, fname), "wb") as f:
                f.write(shape.image.blob)
            from PIL import Image
            import io
            im = Image.open(io.BytesIO(shape.image.blob))
            s_data["images"].append({
                "file": fname,
                "w_px": im.width, "h_px": im.height,
                "x_in": round(Emu(shape.left).inches, 2) if shape.left else 0,
                "y_in": round(Emu(shape.top).inches, 2) if shape.top else 0,
            })
    data.append(s_data)
    print(f"Slide {s_idx:2d}: {len(s_data['texts'])} text, {len(s_data['images'])} image(s)")
    for txt in s_data["texts"]:
        print(f"    text: {txt['text']}")
    for im in s_data["images"]:
        print(f"    img : {im['file']}  ({im['w_px']}×{im['h_px']})")

with open(os.path.join(OUT, "data.json"), "w", encoding="utf-8") as f:
    json.dump(data, f, ensure_ascii=False, indent=2)
print(f"\nSaved to {OUT}")
